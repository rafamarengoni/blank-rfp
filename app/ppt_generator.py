from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import io

import streamlit as st

# Dashboard for presentation customization
st.sidebar.header("Customize Presentation")
theme = st.sidebar.selectbox("Select Theme", ["Corporate", "Modern", "Minimal"])
slide_order = st.sidebar.multiselect("Slide Order", ["Title", "Goals", "Deliverables", "Timeline"], default=["Title", "Goals"])
custom_colors = st.sidebar.color_picker("Pick a Slide Background Color", "#0070C0")


def generate_pptx(details):
    """Generate a PowerPoint presentation from RFP details."""
    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(2))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = "RFP Response Presentation"
    p.font.name = "Helvetica"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Add slides for each detail section
    for section, content in details.items():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(1))
        title_frame = title_box.text_frame
        p = title_frame.add_paragraph()
        p.text = section
        p.font.name = "Helvetica"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(3))
        content_frame = content_box.text_frame
        p = content_frame.add_paragraph()
        p.text = content
        p.font.name = "Helvetica"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
